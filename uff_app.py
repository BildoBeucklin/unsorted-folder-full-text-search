import sys
import os
import sqlite3
import pdfplumber
import numpy as np
import zipfile  
import io       
import traceback

# --- OPTIONALE IMPORTE ---
try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from sentence_transformers import SentenceTransformer, util
from rapidfuzz import process, fuzz

from PyQt6.QtCore import qInstallMessageHandler, QtMsgType, Qt, QThread, pyqtSignal, QUrl, QSize
from PyQt6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, 
                             QHBoxLayout, QLineEdit, QPushButton, QLabel, 
                             QFileDialog, QProgressBar, QMessageBox,
                             QListWidget, QListWidgetItem, QSplitter, QFrame, 
                             QSplashScreen, QScrollArea, QStyle, QGraphicsDropShadowEffect)
from PyQt6.QtGui import QDesktopServices, QPixmap, QCursor, QAction, QColor, QPalette, QFont

# --- 0. LOGGING & SETUP ---

if os.name == 'nt':
    base_dir = os.getenv('LOCALAPPDATA')
else:
    base_dir = os.path.join(os.path.expanduser("~"), ".local", "share")

log_dir = os.path.join(base_dir, "UFF_Search")
if not os.path.exists(log_dir):
    os.makedirs(log_dir)

log_file_path = os.path.join(log_dir, "uff.log")

class Logger(object):
    def __init__(self):
        self.log = open(log_file_path, "w", encoding="utf-8")

    def write(self, message):
        self.log.write(message)
        self.log.flush()

    def flush(self):
        self.log.flush()

sys.stdout = Logger()
sys.stderr = sys.stdout

# --- STYLESHEET ---
STYLESHEET = """
QMainWindow {
    background-color: #f4f7f6;
}

/* Sidebar Styles */
QFrame#Sidebar {
    background-color: #2c3e50;
    border: none;
}
QLabel#SidebarTitle {
    color: #ecf0f1;
    font-weight: bold;
    font-size: 16px;
    padding: 10px;
}
QListWidget {
    background-color: #34495e;
    color: #ecf0f1;
    border: none;
    outline: none;
    font-size: 13px;
}
QListWidget::item {
    padding: 8px;
    border-bottom: 1px solid #2c3e50;
}
QListWidget::item:selected {
    background-color: #1abc9c;
    color: white;
}
QListWidget::item:hover {
    background-color: #16a085;
}

/* Sidebar Buttons */
QPushButton#SidebarBtn {
    background-color: #34495e;
    color: #bdc3c7;
    border: 1px solid #2c3e50;
    padding: 8px;
    text-align: left;
    border-radius: 4px;
    margin: 2px 10px;
}
QPushButton#SidebarBtn:hover {
    background-color: #1abc9c;
    color: white;
    border: 1px solid #16a085;
}
QPushButton#CancelBtn {
    background-color: #e74c3c;
    color: white;
    font-weight: bold;
    border-radius: 4px;
    margin: 10px;
    padding: 8px;
}

/* Main Area */
QLineEdit {
    padding: 10px;
    border: 1px solid #bdc3c7;
    border-radius: 20px;
    font-size: 14px;
    background-color: white;
    selection-background-color: #3498db;
}
QLineEdit:focus {
    border: 2px solid #3498db;
}

QPushButton#SearchBtn {
    background-color: #3498db;
    color: white;
    font-weight: bold;
    border-radius: 20px;
    padding: 10px 20px;
    font-size: 14px;
}
QPushButton#SearchBtn:hover {
    background-color: #2980b9;
}
QPushButton#SearchBtn:pressed {
    background-color: #1f618d;
}

/* Scroll Area & Results */
QScrollArea {
    border: none;
    background-color: transparent;
}
QWidget#ResultsContainer {
    background-color: transparent;
}
QLabel#StatusLabel {
    color: #7f8c8d;
    font-size: 12px;
    margin-left: 10px;
}
QProgressBar {
    border: none;
    background-color: #ecf0f1;
    height: 4px;
    text-align: center;
}
QProgressBar::chunk {
    background-color: #1abc9c;
}
"""

def qt_message_handler(mode, context, message):
    msg_lower = message.lower()
    ignore_keywords = [
        "qt.text.font", "qt.qpa.fonts", "opentype", "directwrite", 
        "unable to create font", "fontbbox", "script"
    ]
    if any(k in msg_lower for k in ignore_keywords): return
    try:
        sys.stdout.write(f"[Qt] {message}\n")
    except: pass 

qInstallMessageHandler(qt_message_handler)
os.environ["QT_LOGGING_RULES"] = "qt.text.font.db=false;qt.qpa.fonts=false"


# --- WIDGET: Modernes Suchergebnis (Fixed Tooltips) ---
class SearchResultItem(QFrame):
    def __init__(self, filename, filepath, snippet, parent=None):
        super().__init__(parent)
        self.filepath = filepath
        
        # WICHTIG: Tooltip auf das gesamte Frame setzen, nicht nur auf Kinder
        self.setToolTip(filepath)
        
        # Design der Karte
        self.setFrameShape(QFrame.Shape.StyledPanel)
        self.setStyleSheet("""
            SearchResultItem {
                background-color: white;
                border: 1px solid #e0e0e0;
                border-radius: 8px;
            }
            SearchResultItem:hover {
                border: 1px solid #3498db;
                background-color: #fbfbfb;
            }
        """)
        
        # Schatten-Effekt
        shadow = QGraphicsDropShadowEffect(self)
        shadow.setBlurRadius(10)
        shadow.setXOffset(0)
        shadow.setYOffset(2)
        shadow.setColor(QColor(0, 0, 0, 30))
        self.setGraphicsEffect(shadow)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(15, 15, 15, 15)
        layout.setSpacing(5)
        
        # 1. Titel (Dateiname)
        self.btn_title = QPushButton(filename)
        self.btn_title.setCursor(Qt.CursorShape.PointingHandCursor)
        # MouseTracking aktivieren hilft manchmal bei schnellen Bewegungen
        self.btn_title.setMouseTracking(True)
        self.btn_title.setStyleSheet("""
            QPushButton {
                text-align: left;
                font-weight: bold;
                font-size: 16px;
                color: #2c3e50;
                border: none;
                background: transparent;
                padding: 0px;
            }
            QPushButton:hover {
                color: #3498db;
                text-decoration: underline;
            }
        """)
        self.btn_title.clicked.connect(self.open_file)
        
        # 2. Snippet
        self.lbl_snippet = QLabel(snippet)
        self.lbl_snippet.setWordWrap(True)
        self.lbl_snippet.setStyleSheet("color: #555; font-size: 13px; line-height: 1.4;")
        
        # 3. Pfad (unten, klein)
        path_layout = QHBoxLayout()
        lbl_icon = QLabel("📄") 
        lbl_icon.setStyleSheet("font-size: 10px; color: #95a5a6;")
        
        self.lbl_path = QLabel(filepath)
        self.lbl_path.setStyleSheet("color: #95a5a6; font-size: 11px;")
        
        path_layout.addWidget(lbl_icon)
        path_layout.addWidget(self.lbl_path)
        path_layout.addStretch()
        
        layout.addWidget(self.btn_title)
        layout.addWidget(self.lbl_snippet)
        layout.addLayout(path_layout)
    
    def open_file(self):
        target_path = self.filepath
        if " :: " in target_path:
            target_path = target_path.split(" :: ")[0]
        url = QUrl.fromLocalFile(target_path)
        QDesktopServices.openUrl(url)


# --- 1. DATENBANK MANAGER ---

class DatabaseHandler:
    def __init__(self):
        self.app_data_dir = log_dir
        self.db_name = os.path.join(self.app_data_dir, "uff_index.db")
        self.model = None 
        self.init_db()

    def init_db(self):
        conn = sqlite3.connect(self.db_name)
        cursor = conn.cursor()
        cursor.execute("CREATE VIRTUAL TABLE IF NOT EXISTS documents USING fts5(filename, path, content);")
        cursor.execute("CREATE TABLE IF NOT EXISTS folders (path TEXT PRIMARY KEY, alias TEXT);")
        cursor.execute("CREATE TABLE IF NOT EXISTS embeddings (doc_id INTEGER PRIMARY KEY, vec BLOB);")
        conn.commit()
        conn.close()

    def add_folder(self, path):
        conn = sqlite3.connect(self.db_name)
        try:
            conn.execute("INSERT OR IGNORE INTO folders (path, alias) VALUES (?, ?)", (path, os.path.basename(path)))
            conn.commit()
            return True
        except: return False
        finally: conn.close()

    def remove_folder(self, path):
        conn = sqlite3.connect(self.db_name)
        cursor = conn.cursor()
        cursor.execute("SELECT rowid FROM documents WHERE path LIKE ?", (f"{path}%",))
        ids = [row[0] for row in cursor.fetchall()]
        if ids:
            cursor.execute("DELETE FROM documents WHERE path LIKE ?", (f"{path}%",))
            cursor.execute(f"DELETE FROM embeddings WHERE doc_id IN ({','.join('?'*len(ids))})", ids)
        cursor.execute("DELETE FROM folders WHERE path = ?", (path,))
        conn.commit()
        conn.close()

    def get_folders(self):
        conn = sqlite3.connect(self.db_name)
        rows = conn.execute("SELECT path FROM folders").fetchall()
        conn.close()
        return [r[0] for r in rows]

    def search(self, query):
        if not query.strip() or not self.model: return []
        
        q_vec = self.model.encode(query, convert_to_tensor=False)
        conn = sqlite3.connect(self.db_name)
        cursor = conn.cursor()
        
        cursor.execute("SELECT doc_id, vec FROM embeddings")
        data = cursor.fetchall()
        doc_ids = [d[0] for d in data]
        if not doc_ids:
            conn.close(); return []

        vecs = np.array([np.frombuffer(d[1], dtype=np.float32) for d in data])
        scores = util.cos_sim(q_vec, vecs)[0].numpy()
        scores = np.clip(scores, 0, 1)
        sem_map = {did: float(s) for did, s in zip(doc_ids, scores)}

        words = query.replace('"', '').split()
        if not words: words = [query]
        fts_query = " OR ".join([f'"{w}"*' for w in words])
        
        try:
            fts_rows = cursor.execute("SELECT rowid, filename, content FROM documents WHERE documents MATCH ? LIMIT 100", (fts_query,)).fetchall()
        except: fts_rows = []

        lex_map = {}
        for did, fname, content in fts_rows:
            r1 = fuzz.partial_ratio(query.lower(), fname.lower())
            r2 = fuzz.partial_token_set_ratio(query.lower(), content[:5000].lower())
            lex_map[did] = max(r1, r2) / 100.0

        final = {}
        ALPHA = 0.65
        BETA = 0.35
        for did, s_score in sem_map.items():
            if s_score < 0.15 and did not in lex_map: continue
            l_score = lex_map.get(did, 0.0)
            h_score = (s_score * ALPHA) + (l_score * BETA)
            if s_score > 0.4 and l_score > 0.6: h_score += 0.1
            final[did] = h_score

        sorted_ids = sorted(final.keys(), key=lambda x: final[x], reverse=True)[:50]
        results = []
        for did in sorted_ids:
            row = cursor.execute("SELECT filename, path, snippet(documents, 2, '<b>', '</b>', '...', 15) FROM documents WHERE rowid = ?", (did,)).fetchone()
            if row: results.append(row)
        conn.close()
        return results

# --- 2. THREADS ---

class ModelLoaderThread(QThread):
    model_loaded = pyqtSignal(object)
    def run(self):
        try:
            model = SentenceTransformer('all-MiniLM-L6-v2')
            self.model_loaded.emit(model)
        except: self.model_loaded.emit(None)

class IndexerThread(QThread):
    progress_signal = pyqtSignal(str)
    finished_signal = pyqtSignal(int, int, bool)

    def __init__(self, folder, db_name, model):
        super().__init__()
        self.folder_path = folder
        self.db_name = db_name
        self.model = model
        self.is_running = True

    def stop(self): self.is_running = False

    def _extract_text(self, stream, filename):
        ext = os.path.splitext(filename)[1].lower()
        text = ""
        try:
            if ext == ".pdf":
                try:
                    with pdfplumber.open(stream) as pdf:
                        for p in pdf.pages:
                            if t := p.extract_text(): text += t + "\n"
                except: pass
            
            elif ext == ".docx" and docx is not None:
                try:
                    doc = docx.Document(stream)
                    for para in doc.paragraphs: text += para.text + "\n"
                except: pass

            elif ext == ".xlsx" and openpyxl is not None:
                try:
                    wb = openpyxl.load_workbook(stream, data_only=True, read_only=True)
                    for sheet in wb.worksheets:
                        text += f"\n--- {sheet.title} ---\n"
                        for row in sheet.iter_rows(values_only=True):
                            row_text = " ".join([str(c) for c in row if c is not None])
                            if row_text.strip(): text += row_text + "\n"
                except: pass

            elif ext == ".pptx" and Presentation is not None:
                try:
                    prs = Presentation(stream)
                    for i, slide in enumerate(prs.slides):
                        text += f"\n--- Folie {i+1} ---\n"
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs: text += run.text + " "
                                    text += "\n"
                except: pass

            elif ext in [".txt", ".md", ".py", ".json", ".csv", ".html", ".log", ".ini", ".xml"]:
                try:
                    content = stream.read()
                    if isinstance(content, str): text = content
                    else: text = content.decode('utf-8', errors='ignore')
                except: pass
        except: pass
        return text

    def run(self):
        conn = sqlite3.connect(self.db_name)
        cursor = conn.cursor()
        
        cursor.execute("SELECT rowid FROM documents WHERE path LIKE ?", (f"{self.folder_path}%",))
        ids = [r[0] for r in cursor.fetchall()]
        if ids:
            cursor.execute("DELETE FROM documents WHERE path LIKE ?", (f"{self.folder_path}%",))
            cursor.execute(f"DELETE FROM embeddings WHERE doc_id IN ({','.join('?'*len(ids))})", ids)
            conn.commit()

        indexed = 0
        skipped = 0
        cancelled = False

        for root, dirs, files in os.walk(self.folder_path):
            if not self.is_running: cancelled = True; break
            for file in files:
                if not self.is_running: cancelled = True; break
                path = os.path.join(root, file)
                self.progress_signal.emit(f"Prüfe: {file}...")

                if file.lower().endswith('.zip'):
                    try:
                        with zipfile.ZipFile(path, 'r') as z:
                            for zi in z.infolist():
                                if zi.is_dir(): continue
                                vpath = f"{path} :: {zi.filename}"
                                with z.open(zi) as zf:
                                    content = self._extract_text(io.BytesIO(zf.read()), zi.filename)
                                    if content and len(content.strip()) > 20:
                                        self._save(cursor, zi.filename, vpath, content)
                                        indexed += 1
                    except: skipped += 1
                else:
                    try:
                        with open(path, "rb") as f:
                            file_content = io.BytesIO(f.read())
                            content = self._extract_text(file_content, file)
                        if content and len(content.strip()) > 20:
                            self._save(cursor, file, path, content)
                            indexed += 1
                        else: skipped += 1
                    except: skipped += 1

            if cancelled: break
        
        conn.commit()
        conn.close()
        self.finished_signal.emit(indexed, skipped, cancelled)

    def _save(self, cursor, fname, path, content):
        cursor.execute("INSERT INTO documents (filename, path, content) VALUES (?, ?, ?)", (fname, path, content))
        did = cursor.lastrowid
        vec = self.model.encode(content[:8000], convert_to_tensor=False).tobytes()
        cursor.execute("INSERT INTO embeddings (doc_id, vec) VALUES (?, ?)", (did, vec))

# --- 3. UI MAIN WINDOW ---

class UffWindow(QMainWindow):
    def __init__(self, splash=None):
        super().__init__()
        self.splash = splash
        self.db = DatabaseHandler()
        self.initUI()
        self.load_saved_folders()

    def initUI(self):
        self.setWindowTitle("UFF Search v7.2 (Stable Tooltips)")
        self.resize(1100, 750)
        
        self.setStyleSheet(STYLESHEET)
        
        central = QWidget()
        self.setCentralWidget(central)
        main_layout = QHBoxLayout(central)
        main_layout.setContentsMargins(0, 0, 0, 0)
        main_layout.setSpacing(0)

        # -- SIDEBAR --
        left_panel = QFrame()
        left_panel.setObjectName("Sidebar")
        left_panel.setFixedWidth(260)
        left_layout = QVBoxLayout(left_panel)
        left_layout.setContentsMargins(0, 20, 0, 20)
        
        lbl_title = QLabel(" UFF SEARCH")
        lbl_title.setObjectName("SidebarTitle")
        
        self.folder_list = QListWidget()
        
        icon_add = self.style().standardIcon(QStyle.StandardPixmap.SP_FileDialogNewFolder)
        icon_del = self.style().standardIcon(QStyle.StandardPixmap.SP_TrashIcon)
        icon_refresh = self.style().standardIcon(QStyle.StandardPixmap.SP_BrowserReload)
        icon_stop = self.style().standardIcon(QStyle.StandardPixmap.SP_DialogCancelButton)

        btn_add = QPushButton(" Ordner hinzufügen")
        btn_add.setObjectName("SidebarBtn")
        btn_add.setIcon(icon_add)
        btn_add.clicked.connect(self.add_new_folder)
        
        btn_del = QPushButton(" Ordner entfernen")
        btn_del.setObjectName("SidebarBtn")
        btn_del.setIcon(icon_del)
        btn_del.clicked.connect(self.delete_selected_folder)
        
        self.btn_rescan = QPushButton(" Neu scannen")
        self.btn_rescan.setObjectName("SidebarBtn")
        self.btn_rescan.setIcon(icon_refresh)
        self.btn_rescan.clicked.connect(self.rescan_selected_folder)
        
        self.btn_cancel = QPushButton("STOPPEN")
        self.btn_cancel.setObjectName("CancelBtn")
        self.btn_cancel.setIcon(icon_stop)
        self.btn_cancel.clicked.connect(self.cancel_indexing)
        self.btn_cancel.hide()

        left_layout.addWidget(lbl_title)
        left_layout.addSpacing(10)
        left_layout.addWidget(self.folder_list)
        left_layout.addSpacing(10)
        left_layout.addWidget(btn_add)
        left_layout.addWidget(btn_del)
        left_layout.addWidget(self.btn_rescan)
        left_layout.addWidget(self.btn_cancel)

        # -- RECHTS (Hauptbereich) --
        right_panel = QWidget()
        right_panel.setObjectName("MainArea")
        right_layout = QVBoxLayout(right_panel)
        right_layout.setContentsMargins(30, 30, 30, 30)
        right_layout.setSpacing(15)
        
        # Header
        search_box = QHBoxLayout()
        self.input_search = QLineEdit()
        self.input_search.setPlaceholderText("Wonach suchst du heute?")
        self.input_search.returnPressed.connect(self.perform_search)
        
        self.btn_go = QPushButton("Suchen")
        self.btn_go.setObjectName("SearchBtn")
        self.btn_go.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_go.clicked.connect(self.perform_search)
        
        search_box.addWidget(self.input_search)
        search_box.addWidget(self.btn_go)

        # Status
        status_box = QHBoxLayout()
        self.lbl_status = QLabel("Modell wird geladen...")
        self.lbl_status.setObjectName("StatusLabel")
        self.progress_bar = QProgressBar()
        self.progress_bar.hide()
        status_box.addWidget(self.lbl_status)
        status_box.addWidget(self.progress_bar)

        # Ergebnisse
        self.scroll_area = QScrollArea()
        self.scroll_area.setWidgetResizable(True)
        
        self.results_container = QWidget()
        self.results_container.setObjectName("ResultsContainer")
        self.results_layout = QVBoxLayout(self.results_container)
        self.results_layout.setAlignment(Qt.AlignmentFlag.AlignTop)
        self.results_layout.setSpacing(15)
        self.scroll_area.setWidget(self.results_container)

        right_layout.addLayout(search_box)
        right_layout.addLayout(status_box)
        right_layout.addWidget(self.scroll_area)

        main_layout.addWidget(left_panel)
        main_layout.addWidget(right_panel)
        self.set_ui_enabled(False)

    def set_ui_enabled(self, enabled):
        self.input_search.setEnabled(enabled)
        self.btn_go.setEnabled(enabled)
        self.folder_list.setEnabled(enabled)

    def start_model_loading(self):
        if self.splash: self.splash.showMessage("Lade KI-Modell...", Qt.AlignmentFlag.AlignBottom, Qt.GlobalColor.white)
        self.loader = ModelLoaderThread()
        self.loader.model_loaded.connect(self.on_model_loaded)
        self.loader.start()

    def on_model_loaded(self, model):
        if self.splash: self.splash.finish(self)
        if not model:
            QMessageBox.critical(self, "Fehler", "Modell konnte nicht geladen werden.")
            return
        self.db.model = model
        self.lbl_status.setText("Bereit für deine Suche.")
        self.set_ui_enabled(True)

    def perform_search(self):
        query = self.input_search.text()
        if not query: return
        self.lbl_status.setText("Suche läuft...")
        QApplication.processEvents() 

        while self.results_layout.count():
            child = self.results_layout.takeAt(0)
            if child.widget(): child.widget().deleteLater()

        results = self.db.search(query)
        self.lbl_status.setText(f"{len(results)} Treffer gefunden.")

        if not results:
            lbl = QLabel("Leider keine Ergebnisse.")
            lbl.setStyleSheet("color: #95a5a6; font-size: 18px; margin-top: 40px;")
            lbl.setAlignment(Qt.AlignmentFlag.AlignHCenter)
            self.results_layout.addWidget(lbl)
        else:
            for fname, fpath, snippet in results:
                self.results_layout.addWidget(SearchResultItem(fname, fpath, snippet))
        self.results_layout.addStretch()

    def load_saved_folders(self):
        self.folder_list.clear()
        for f in self.db.get_folders():
            item = QListWidgetItem(self.style().standardIcon(QStyle.StandardPixmap.SP_DirIcon), f)
            item.setToolTip(f)
            self.folder_list.addItem(item)

    def add_new_folder(self):
        f = QFileDialog.getExistingDirectory(self, "Ordner wählen")
        if f and self.db.add_folder(f):
            self.load_saved_folders()
            self.start_idx(f)

    def delete_selected_folder(self):
        item = self.folder_list.currentItem()
        if item and QMessageBox.question(self, "Löschen", f"Weg damit?\n{item.text()}", QMessageBox.StandardButton.Yes|QMessageBox.StandardButton.No) == QMessageBox.StandardButton.Yes:
            self.db.remove_folder(item.text())
            self.load_saved_folders()

    def rescan_selected_folder(self):
        if item := self.folder_list.currentItem(): self.start_idx(item.text())

    def start_idx(self, folder):
        if not self.db.model: return
        self.set_ui_enabled(False)
        self.btn_cancel.show(); self.btn_rescan.hide(); self.progress_bar.show()
        self.idx_thread = IndexerThread(folder, self.db.db_name, self.db.model)
        self.idx_thread.progress_signal.connect(self.lbl_status.setText)
        self.idx_thread.finished_signal.connect(self.idx_done)
        self.idx_thread.start()

    def cancel_indexing(self):
        if self.idx_thread: self.idx_thread.stop()

    def idx_done(self, n, s, c):
        self.set_ui_enabled(True)
        self.btn_cancel.hide(); self.btn_rescan.show(); self.progress_bar.hide()
        msg = "Abgebrochen" if c else "Indexierung fertig"
        self.lbl_status.setText(f"{msg}: {n} neu, {s} übersprungen.")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    
    font = QFont("Segoe UI", 10)
    app.setFont(font)

    splash = None
    try:
        if os.path.exists("assets/uff_banner.jpeg"):
            splash = QSplashScreen(QPixmap("assets/uff_banner.jpeg"))
            splash.show()
    except: pass
    w = UffWindow(splash)
    w.show()
    w.start_model_loading()
    sys.exit(app.exec())